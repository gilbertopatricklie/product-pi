import os
from flask import Flask, request, render_template
import pymupdf
from pptx import Presentation
from gtts import gTTS
from transformers import T5Tokenizer, T5ForConditionalGeneration
import torch
import re

#init flask
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['AUDIO_FOLDER'] = 'static'

#load model
model_name = "cahya/t5-base-indonesian-summarization-cased"
tokenizer = T5Tokenizer.from_pretrained(model_name)
model = T5ForConditionalGeneration.from_pretrained(model_name)

#bersihin teks
def clean_text(text):
    text = re.sub(r'[•●▪◦]', '', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

#ekstrak teks dari pdf
def extract_text_from_pdf(filepath):
    doc = pymupdf.open(filepath)
    raw_text = "\n".join(page.get_text() for page in doc)
    return clean_text(raw_text)

#ekstrak teks dari pdf
def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return clean_text(text)

def chunk_text(text, max_words=200, overlap=50):
    words = text.split()
    chunks = []
    i = 0
    
    while i < len(words):
        chunk = words[i:i + max_words]
        chunks.append(" ".join(chunk))
        i += max_words - overlap
    return chunks

def summarize_chunk(chunk_text, max_output=150):
    input_text = "Hasil Ringkasan: " + chunk_text
    inputs = tokenizer.encode(input_text, return_tensors="pt", truncation=True, max_length=512)
    with torch.no_grad():
        summary_ids = model.generate(
            inputs,
            max_length=max_output,
            min_length=30,
            num_beams=4,
            length_penalty=2.0,
            early_stopping=True
        )
    return tokenizer.decode(summary_ids[0], skip_special_tokens=True)

def summarize_text(text, max_output=150):
   chunks = chunk_text(text)
   summaries = [summarize_chunk(chunk, max_output=max_output) for chunk in chunks]
   final_summary = " ".join(summaries)
   return final_summary

def text_to_audio(text, output_path):
    tts = gTTS(text, lang='id')
    tts.save(output_path)

@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    uploaded_file = request.files['file']
    output_name = request.form.get('output_name', 'ringkasan_audio').strip()
    output_name = "".join(c for c in output_name if c.isalnum() or c in (' ', '_', '-')).rstrip()
    filename = uploaded_file.filename

    if not uploaded_file or not filename.endswith(('.pdf', '.pptx')):
        return "Format file tidak didukung.", 400

    # Simpan file sementara
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    uploaded_file.save(filepath)

    # Ekstraksi teks
    if filename.endswith('.pdf'):
        text = extract_text_from_pdf(filepath)
    else:
        text = extract_text_from_pptx(filepath)

    # Summarize dan konversi ke audio
    summary = summarize_text(text)
    audio_filename = output_name + ".mp3"
    audio_path = os.path.join(app.config['AUDIO_FOLDER'], audio_filename)
    text_to_audio(summary, audio_path)
    
    # Simpan ringkasan ke file .txt
    summary_filename = output_name + ".txt"
    summary_path = os.path.join(app.config['AUDIO_FOLDER'], summary_filename)
    with open(summary_path, "w", encoding="utf-8") as f:
        f.write(summary)
        
    audio_url = f"/static/{audio_filename}"
    summary_url = f"/static/{summary_filename}"
    
    return render_template(
    'index.html',
    summary=summary,
    summary_url=summary_url,
    audio_url=audio_url
)

if __name__ == '__main__':
    os.makedirs('uploads', exist_ok=True)
    os.makedirs('static', exist_ok=True)
    app.run(debug=True)