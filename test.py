import re
import pymupdf
from pptx import Presentation

##BUAT TESTING FUNGSI AMA NUNJUKIN HASIL OUTPUT DOANG

filepath = "cocomo.pptx"

def clean_text(text):
    text = re.sub(r'[•●▪◦]', '', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return clean_text(text)

def chunk_text(text, max_words=200, overlap=50):
    words = text.split()
    chunks = []
    i = 0
    
    while i < len(words):
        chunk = words[i:i + max_words]
        chunks.append(" ".join(chunk))
        i += max_words - overlap
    return chunks

text = extract_text_from_pptx(filepath)
chunks = chunk_text(text, max_words=200, overlap=50)

for i, chunk in enumerate(chunks):
    print(f"\n")
    print(chunk)